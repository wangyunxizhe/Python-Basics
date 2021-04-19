import pptx
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION

# 步骤1：得到演示文稿对象。若是对现有的ppt文件进行改动，在入参中填入需要改动的ppt地址即可；若新键ppt则空参
prs = pptx.Presentation()

# 步骤2：写入操作
# 设置ppt的幻灯片布局
slide = prs.slides.add_slide(prs.slide_layouts[0])  # 第一张幻灯片用 样式1
prs.slides.add_slide(prs.slide_layouts[1])  # 第二张幻灯片用 样式2
prs.slides.add_slide(prs.slide_layouts[2])  # 第三张幻灯片用 样式3
print("幻灯片的总数：", len(prs.slides))
# 删除幻灯片
del prs.slides._sldIdLst[1]
print("删除后，幻灯片的总数：", len(prs.slides))
text1 = slide.shapes.add_textbox(Inches(5), Inches(5), Inches(5), Inches(5))  # 距离上下左右各5英寸的距离
text1.text = "我是文本框"
p1 = text1.text_frame.add_paragraph()
p1.text = "我是段落1"
p1.add_run().text = "end"
title_shape = slide.shapes.title
title_shape.text = "标题1"
slide.shapes.placeholders[1].text = "标题2"
# 添加自选图形
shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2), Inches(2), Inches(5), Inches(3))
# 填充，边框
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)
line = shape.line
line.color.rgb = RGBColor(55, 3, 5)
line.width = Pt(2)
# 添加表格
table = slide.shapes.add_table(3, 3, Inches(2), Inches(2), Inches(4), Inches(2)).table  # 两行三列
# 填充内容
table.cell(1, 0).text = "name"
table.cell(1, 1).text = "age"
table.cell(1, 2).text = "class"
table.cell(2, 0).text = "张三"
table.cell(2, 1).text = "19"
table.cell(2, 2).text = "一班"
# 合并单元格
cell = table.cell(0, 0)
cell1 = table.cell(0, 2)
cell.merge(cell1)
table.cell(0, 0).text = "班级学生信息"  # 第一行
print("单元格是否合并：", cell.is_merge_origin)
# 取消合并
# cell.split()
# 写入图表
chart_data = CategoryChartData()
chart_data.categories = ["一月份", "二月份", "三月份"]  # X轴
chart_data.add_series("Y2020", (300, 400, 500))
chart_data.add_series("Y2021", (400, 500, 600))
chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(2), Inches(2), Inches(6), Inches(4), chart_data).chart
chart.has_title = True
chart.chart_title.text_frame.text = "第一季度销售额"  # 标题
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT

# 步骤3：保存PPT文件
prs.save("../../resources/test.pptx")
